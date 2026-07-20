import logging
import re
import time
from datetime import datetime
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from ..config import DEFAULT_REPORT_AUDIENCE, data_dir
from ..models import AppSetting, Evidence, EvidenceRelationship, GeneratedReport, ReportRevision, ReportStyleProfile, WorkItem, Workspace
from .ai import AIUnavailable, OllamaProvider
from .report_safety import build_report_plan, plan_to_markdown, presentation_plan_from_report_plan, refine_markdown, validate_and_rewrite

logger = logging.getLogger(__name__)


REPORT_TYPES = [
    "Daily Work Summary",
    "Weekly Work Report",
    "Monthly Work Summary",
    "Stakeholder Update",
    "Issue and Resolution Report",
    "Investigation Report",
    "Project Progress Report",
    "Meeting Brief",
    "Handover Document",
    "CV Achievement Extraction",
    "Meeting Presentation",
]

LEGACY_REPORT_TYPE_ALIASES = {
    "Monthly Engineering Summary": "Monthly Work Summary",
    "Boss Progress Update": "Stakeholder Update",
    "Bug and Resolution Report": "Issue and Resolution Report",
    "Technical Investigation Report": "Investigation Report",
}

ACCEPTED_REPORT_TYPES = set(REPORT_TYPES) | set(LEGACY_REPORT_TYPE_ALIASES)


def evidence_for_period(db: Session, workspace_id: int, date_from: str, date_to: str, project_id: int | None = None):
    promoted_source_item_ids = promoted_work_item_ids(db)
    confirmed_query = (
        db.query(WorkItem)
        .filter(WorkItem.workspace_id == workspace_id)
        .filter(WorkItem.work_date >= date_from, WorkItem.work_date <= date_to)
        .filter(WorkItem.status == "CONFIRMED")
    )
    if project_id is not None:
        confirmed_query = confirmed_query.filter(WorkItem.project_id == project_id)
    confirmed = confirmed_query.order_by(WorkItem.work_date).all()
    confirmed = [item for item in confirmed if item.id not in promoted_source_item_ids]
    inferred_query = (
        db.query(WorkItem)
        .filter(WorkItem.workspace_id == workspace_id)
        .filter(WorkItem.work_date >= date_from, WorkItem.work_date <= date_to)
        .filter(WorkItem.status == "REVIEW")
    )
    if project_id is not None:
        inferred_query = inferred_query.filter(WorkItem.project_id == project_id)
    inferred = inferred_query.order_by(WorkItem.work_date).all()
    attach_evidence_summaries(db, confirmed + inferred)
    return confirmed, inferred


def promoted_work_item_ids(db: Session) -> set[int]:
    relationships = db.query(EvidenceRelationship).filter(EvidenceRelationship.relationship_type == "PROMOTED_TO").all()
    source_ids = [row.from_evidence_id for row in relationships]
    if not source_ids:
        return set()
    evidence_rows = db.query(Evidence).filter(Evidence.id.in_(source_ids), Evidence.work_item_id.isnot(None)).all()
    return {row.work_item_id for row in evidence_rows if row.work_item_id}


def generate_report(
    db: Session,
    workspace_id: int,
    report_type: str,
    date_from: str,
    date_to: str,
    include_inferred_ids: list[int] | None = None,
    project_id: int | None = None,
) -> GeneratedReport:
    workspace = db.get(Workspace, workspace_id)
    workspace_name = workspace.name if workspace else "Current Workspace"
    confirmed, _ = evidence_for_period(db, workspace_id, date_from, date_to, project_id)
    included = list(confirmed)
    if include_inferred_ids:
        inferred_query = db.query(WorkItem).filter(WorkItem.id.in_(include_inferred_ids), WorkItem.workspace_id == workspace_id)
        if project_id is not None:
            inferred_query = inferred_query.filter(WorkItem.project_id == project_id)
        included.extend(inferred_query.all())
    profile = db.query(ReportStyleProfile).filter(ReportStyleProfile.workspace_id == workspace_id).first()
    display_type = report_type_label(report_type)
    plan = build_report_plan(display_type, date_from, date_to, included)
    generated = ai_report(db, display_type, date_from, date_to, included, profile, workspace_name) if included else None
    analysis_mode = "OLLAMA" if generated else "EVIDENCE_ONLY"
    markdown, notes = validate_and_rewrite(generated or deterministic_report(display_type, date_from, date_to, included, profile, workspace_name), plan)
    report = GeneratedReport(
        workspace_id=workspace_id,
        project_id=project_id,
        report_type=report_type,
        title=f"{display_type} - {date_from} to {date_to}",
        date_from=date_from,
        date_to=date_to,
        draft_markdown=markdown,
    )
    report.analysis_mode = analysis_mode
    db.add(report)
    db.flush()
    db.add(ReportRevision(report_id=report.id, revision_number=1, reason="Initial generation", draft_markdown=markdown, validation_notes="\n".join(notes)))
    db.commit()
    db.refresh(report)
    return report


def ai_report(
    db: Session,
    report_type: str,
    date_from: str,
    date_to: str,
    items: list[WorkItem],
    profile: ReportStyleProfile | None,
    workspace_name: str,
) -> str | None:
    model_setting = db.get(AppSetting, "selected_model")
    model = model_setting.value if model_setting else ""
    if not model:
        return None
    provider = OllamaProvider()
    if not provider.health_check().get("available"):
        return None
    evidence = "\n".join(f"- {item.work_date}: {item.title} ({item.work_type}; {item.work_status or 'status unknown'}) - {item.summary}{general_item_suffix(item)}" for item in items)
    audience = profile.audience if profile and profile.audience else DEFAULT_REPORT_AUDIENCE
    prompt = f"""
Create a grounded {report_type} for the workspace named {workspace_name}.
Period: {date_from} to {date_to}
Audience: {audience}
Tone: {profile.tone if profile else 'Professional'}
Rules:
- Use only the evidence below.
- Do not invent accomplishments.
- Do not include empty sections.
- Mark pending work clearly.

    Evidence:
{evidence or 'No confirmed work items are available.'}
"""
    start = time.perf_counter()
    try:
        text = provider.generate_text(model, prompt, system="You write accurate professional work reports from local evidence only.")
        logger.info(
            "ai_report mode=OLLAMA provider=Ollama model=%s report_type=%s duration=%.3fs items=%s",
            model,
            report_type,
            round(time.perf_counter() - start, 3),
            len(items),
        )
        return text
    except AIUnavailable as exc:
        logger.warning("ai_report fallback=deterministic report_type=%s reason=%s", report_type, exc)
        return None


def deterministic_report(
    report_type: str,
    date_from: str,
    date_to: str,
    items: list[WorkItem],
    profile: ReportStyleProfile | None,
    workspace_name: str,
) -> str:
    plan = build_report_plan(report_type, date_from, date_to, items)
    if report_type == "Stakeholder Update":
        return deterministic_stakeholder_update(date_from, date_to, items, workspace_name)
    if report_type == "CV Achievement Extraction":
        return deterministic_cv_achievements(date_from, date_to, items, workspace_name)
    return plan_to_markdown(plan)


def deterministic_stakeholder_update(date_from: str, date_to: str, items: list[WorkItem], workspace_name: str) -> str:
    lines = ["# Stakeholder Update", "", f"Period: {date_from} to {date_to}", ""]
    if not items:
        return "\n".join(lines + ["No confirmed work evidence is available for this period yet."])
    lines.append(f"Here is a concise update on confirmed work for {workspace_name}:")
    lines.append("")
    for item in items[:6]:
        lines.append(f"- {item.summary}{general_item_suffix(item)}")
    pending = [item.next_step or item.pending_work for item in items if item.next_step or item.pending_work]
    if pending:
        lines.extend(["", "Pending / next:", *[f"- {text}" for text in pending[:3]]])
    return "\n".join(lines)


def deterministic_cv_achievements(date_from: str, date_to: str, items: list[WorkItem], workspace_name: str) -> str:
    lines = ["# CV Achievement Extraction", "", f"Evidence period: {date_from} to {date_to}", ""]
    if not items:
        return "\n".join(lines + ["No confirmed work evidence is available for CV achievement extraction."])
    by_area: dict[str, list[WorkItem]] = {}
    for item in items:
        by_area.setdefault(item.area or "General Work", []).append(item)
    for area, group in by_area.items():
        lines.append(f"## {area}")
        lines.append(
            f"- Suggested CV bullet: Contributed to {area.lower()} work in {workspace_name} through {len(group)} confirmed evidence-backed activity record(s)."
        )
        lines.append(f"- Supporting evidence: {', '.join(item.title for item in group[:5])}")
        lines.append("- Confidence: CONFIRMED evidence, wording requires user review before use.")
        lines.append("")
    return "\n".join(lines)


def export_docx(report: GeneratedReport) -> Path:
    doc = Document()
    if report.status != "APPROVED":
        doc.add_paragraph("DRAFT - not approved for final use.")
    for line in report.draft_markdown.splitlines():
        stripped = line.strip()
        if not stripped or stripped == "---":
            continue
        if internal_metadata(remove_markdown_prefix(stripped)):
            continue
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:].strip(), level=0)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:].strip(), level=1)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:].strip(), level=2)
        elif re.match(r"^\d+\.\s+", stripped):
            add_markdown_runs(doc.add_paragraph(style="List Number"), re.sub(r"^\d+\.\s+", "", stripped))
        elif stripped.startswith("- "):
            add_markdown_runs(doc.add_paragraph(style="List Bullet"), stripped[2:].strip())
        elif not internal_metadata(stripped):
            add_markdown_runs(doc.add_paragraph(), stripped)
    path = data_dir() / "exports" / f"report-{report.id}.docx"
    doc.save(path)
    return path


def export_pptx(report: GeneratedReport, items: list[WorkItem] | None = None) -> Path:
    prs = Presentation()
    if items is not None:
        report_plan = build_report_plan(report.report_type, report.date_from, report.date_to, items)
        plan = presentation_plan_from_report_plan(report_plan)
    else:
        plan = presentation_plan_from_report(report)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = plan["title"]
    title_slide.placeholders[1].text = plan["reporting_period"] if report.status == "APPROVED" else f"{plan['reporting_period']} - DRAFT"
    for slide_data in plan["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"][:80]
        body = slide.placeholders[1].text_frame
        body.clear()
        for bullet in slide_data["bullets"][:5]:
            p = body.add_paragraph()
            p.text = bullet[:140]
            p.font.size = Pt(18)
    path = data_dir() / "exports" / f"report-{report.id}.pptx"
    prs.save(path)
    return path


def parse_sections(markdown: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current = None
    bullets: list[str] = []
    for line in markdown.splitlines():
        if line.startswith("## "):
            if current:
                sections.append((current, bullets))
            current = line[3:].strip()
            bullets = []
        elif line.startswith("- "):
            bullets.append(line[2:].replace("**", ""))
        elif current and line.strip() and len(bullets) < 5:
            bullets.append(line.strip().replace("#", ""))
    if current:
        sections.append((current, bullets))
    return sections or [("Progress Overview", [line.strip("# ") for line in markdown.splitlines() if line.strip()][:5])]


def create_report_revision(db: Session, report: GeneratedReport, instruction: str) -> ReportRevision:
    items = report_items_for_revision(db, report)
    plan = build_report_plan(report.report_type, report.date_from, report.date_to, items)
    revised, notes = refine_markdown(report.draft_markdown, instruction, plan)
    max_revision = db.query(ReportRevision).filter(ReportRevision.report_id == report.id).order_by(ReportRevision.revision_number.desc()).first()
    next_number = (max_revision.revision_number if max_revision else 0) + 1
    revision = ReportRevision(
        report_id=report.id,
        revision_number=next_number,
        reason=instruction[:260],
        draft_markdown=revised,
        validation_notes="\n".join(notes),
    )
    report.draft_markdown = revised
    report.status = "DRAFT"
    report.approved_at = None
    db.add(revision)
    db.commit()
    db.refresh(revision)
    return revision


def restore_report_revision(db: Session, report: GeneratedReport, revision: ReportRevision):
    report.draft_markdown = revision.draft_markdown
    report.status = "DRAFT"
    report.approved_at = None
    db.commit()
    return report


def report_items_for_revision(db: Session, report: GeneratedReport) -> list[WorkItem]:
    confirmed, _ = evidence_for_period(db, report.workspace_id, report.date_from, report.date_to, report.project_id)
    return confirmed


def attach_evidence_summaries(db: Session, items: list[WorkItem]):
    item_ids = [item.id for item in items]
    if not item_ids:
        return
    rows = (
        db.query(Evidence)
        .filter(Evidence.work_item_id.in_(item_ids), Evidence.archived_at.is_(None))
        .order_by(Evidence.created_at.desc())
        .all()
    )
    by_item: dict[int, list[Evidence]] = {}
    for row in rows:
        if row.work_item_id:
            by_item.setdefault(row.work_item_id, []).append(row)
    for item in items:
        item.evidence_summaries = [
            f"{row.evidence_type or row.source_type}: {row.title} - {row.summary}".strip(" -")
            for row in by_item.get(item.id, [])[:5]
        ]


def general_item_suffix(item: WorkItem) -> str:
    parts = []
    if item.category:
        parts.append(f"Category: {item.category}")
    if item.priority and item.priority != "NORMAL":
        parts.append(f"Priority: {item.priority}")
    if item.outcome:
        parts.append(f"Outcome: {item.outcome}")
    if item.next_step:
        parts.append(f"Next step: {item.next_step}")
    evidence_summaries = getattr(item, "evidence_summaries", [])
    if evidence_summaries:
        parts.append("Evidence: " + " | ".join(evidence_summaries[:3]))
    return (" " + " ".join(parts)) if parts else ""


def add_markdown_runs(paragraph, text: str):
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        else:
            paragraph.add_run(part)


def internal_metadata(text: str) -> bool:
    return bool(re.match(r"^(Audience|Tone|Final Date):", remove_markdown_prefix(text), flags=re.I))


def remove_markdown_prefix(text: str) -> str:
    return re.sub(r"^#{1,6}\s*", "", text or "").strip()


def presentation_plan_from_report(report: GeneratedReport) -> dict:
    sections = parse_sections(report.draft_markdown)
    slides = []
    desired = {
        "Executive Summary": "Progress Overview",
        "Context": "Context",
        "Work Completed / Investigated": "Work Investigated / Completed",
        "Confirmed Findings": "Key Findings",
        "Pending Work": "Open Questions / Pending Work",
        "Next Steps": "Next Steps",
    }
    for heading, bullets in sections:
        title = desired.get(heading, heading)
        clean_bullets = [clean_slide_bullet(bullet) for bullet in bullets if clean_slide_bullet(bullet)]
        if clean_bullets:
            slides.append({"title": title, "purpose": title, "bullets": clean_bullets[:5], "evidence_refs": []})
    if not slides:
        slides.append({"title": "Progress Overview", "purpose": "overview", "bullets": ["No confirmed evidence is available for this period."], "evidence_refs": []})
    return {
        "title": "Meeting Presentation" if report_type_label(report.report_type) == "Meeting Presentation" else report.title,
        "reporting_period": f"{report.date_from} to {report.date_to}",
        "slides": slides[:6],
    }


def clean_slide_bullet(text: str) -> str:
    text = re.sub(r"\*\*", "", text)
    text = re.sub(r"^#+\s*", "", text)
    text = re.sub(r"^(Audience|Tone|Final Date):.*$", "", text, flags=re.I)
    text = text.strip(" -")
    if not text or text == "---":
        return ""
    return text


def report_type_label(report_type: str) -> str:
    return LEGACY_REPORT_TYPE_ALIASES.get(report_type, report_type)
