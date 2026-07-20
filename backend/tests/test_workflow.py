import os
import shutil
import subprocess
import sys
from pathlib import Path

from docx import Document
from pptx import Presentation

BACKEND_ROOT = Path(__file__).resolve().parents[1]
if str(BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(BACKEND_ROOT))

from app.services.report_safety import build_report_plan, normalized_semantic_text, presentation_plan_from_report_plan
from app.services.ai import OllamaProvider
from app.services.semantic_extraction import RawSemanticExtractionSchema, SemanticEventSchema, SemanticExtractionSchema, SemanticFactSchema, normalize_semantic_result


NOTIFICATION_NOTE = (
    "okay so today I was checking push notifications because boss asked if we can activate company notifications "
    "from the dashboard. initially I thought he meant ride and boat cruise notifications but he clarified he means "
    "things like promos and price drops. I checked the dashboard configurations and found notification toggles but "
    "I still need to confirm whether there's an actual broadcast notification flow"
)

MULTI_EVENT_NOTE = (
    "today i fixed the boat cruise category images because one of the icons was wrong then later i checked notifications "
    "after boss clarified he wants company announcements not booking notifications. i found notification toggles in the "
    "dashboard but i haven't confirmed if broadcast sending is actually implemented"
)


def test_work_log_creates_review_item_and_confirm(client):
    response = client.post(
        "/api/work-logs",
        json={"raw_text": "Checked dashboard notification config and investigated company promo notification workflow."},
    )
    assert response.status_code == 200
    items = response.json()["extracted_items"]
    assert items
    assert items[0]["status"] == "REVIEW"

    confirm = client.post(f"/api/work-items/{items[0]['id']}/confirm")
    assert confirm.status_code == 200
    assert confirm.json()["status"] == "CONFIRMED"
    assert confirm.json()["evidence_confidence"] == "CONFIRMED"


def test_default_workspace_is_preserved_and_startup_is_idempotent(client):
    first = client.get("/api/workspaces/default").json()
    second = client.get("/api/workspaces/default").json()
    assert first["id"] == second["id"]
    assert first["name"] == "Ride Yanga"

    from app.bootstrap import ensure_defaults
    from app.db import SessionLocal
    from app.models import Workspace

    db = SessionLocal()
    try:
        ensure_defaults(db)
        ensure_defaults(db)
        rows = db.query(Workspace).all()
        assert len(rows) == 1
        assert rows[0].name == "Ride Yanga"
    finally:
        db.close()


def test_workspace_rename_does_not_create_duplicate_default_workspace(client):
    client.put("/api/workspaces/default", json={"name": "Freelance Work", "user_name": "Alfy"})
    assert client.get("/api/workspaces/default").json()["name"] == "Freelance Work"

    from app.bootstrap import ensure_defaults
    from app.db import SessionLocal
    from app.models import Workspace

    db = SessionLocal()
    try:
        ensure_defaults(db)
        rows = db.query(Workspace).all()
        assert len(rows) == 1
        assert rows[0].name == "Freelance Work"
    finally:
        db.close()


def test_user_display_name_and_report_audience_are_configurable(client):
    client.put("/api/settings", json={"key": "profile_display_name", "value": "Miriam"})
    response = client.put(
        "/api/workspaces/default",
        json={
            "name": "Freelance Work",
            "role_title": "Designer",
            "report_audience": "Client",
        },
    )
    assert response.status_code == 200
    assert response.json()["name"] == "Freelance Work"

    settings = client.get("/api/settings").json()
    assert settings["settings"]["profile_display_name"] == "Miriam"
    assert settings["settings"]["profile_role_title"] == "Designer"
    assert settings["settings"]["default_report_audience"] == "Client"
    assert settings["style_profile"]["audience"] == "Client"


def test_workspace_rename_does_not_change_profile_identity(client):
    client.put("/api/settings", json={"key": "profile_display_name", "value": "Miriam"})
    client.put("/api/workspaces/default", json={"name": "Academic Work"})
    settings = client.get("/api/settings").json()
    assert settings["settings"]["profile_display_name"] == "Miriam"
    assert client.get("/api/workspaces/default").json()["name"] == "Academic Work"


def test_multiple_workspaces_active_selection_and_fallback(client):
    ride = client.get("/api/workspaces/active").json()
    created = client.post("/api/workspaces", json={"name": "Freelance Work", "workspace_type": "Freelance"})
    assert created.status_code == 200
    assert len(client.get("/api/workspaces").json()) == 2
    duplicate = client.post("/api/workspaces", json={"name": "Freelance Work"})
    assert duplicate.status_code == 400

    selected = client.post(f"/api/workspaces/{created.json()['id']}/select")
    assert selected.status_code == 200
    assert client.get("/api/workspaces/active").json()["name"] == "Freelance Work"

    from app.bootstrap import ensure_defaults
    from app.db import SessionLocal
    from app.models import AppSetting

    db = SessionLocal()
    try:
        db.get(AppSetting, "active_workspace_id").value = "999999"
        db.commit()
        fallback = ensure_defaults(db)
        assert fallback.id == ride["id"]
        assert db.get(AppSetting, "active_workspace_id").value == str(ride["id"])
    finally:
        db.close()


def test_switching_workspace_scopes_work_items_search_and_chat(client):
    ride_item = client.post("/api/work-logs", json={"raw_text": "Completed Ride Yanga booking follow-up."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{ride_item['id']}/confirm")
    freelance = client.post("/api/workspaces", json={"name": "Freelance Work"}).json()
    client.post(f"/api/workspaces/{freelance['id']}/select")
    freelance_item = client.post("/api/work-logs", json={"raw_text": "Met a client about a poster revision."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{freelance_item['id']}/confirm")

    items = client.get("/api/work-items").json()
    assert {item["id"] for item in items} == {freelance_item["id"]}
    assert not client.get("/api/search?q=booking").json()
    assert client.get("/api/search?q=poster").json()

    chat = client.post("/api/chat", json={"message": "What work did I complete this week?"}).json()
    assert "poster" in chat["answer"].lower()
    assert "booking" not in chat["answer"].lower()


def test_projects_are_scoped_validated_updated_and_archived(client):
    missing = client.post("/api/projects", json={"workspace_id": 999999, "name": "Missing Workspace Project"})
    assert missing.status_code == 404

    first = client.post("/api/projects", json={"name": "Client Website", "category": "Client", "description": "Website work"})
    assert first.status_code == 200
    project = first.json()
    duplicate = client.post("/api/projects", json={"name": "Client Website"})
    assert duplicate.status_code == 400

    updated = client.put(
        f"/api/projects/{project['id']}",
        json={"name": "Client Website Refresh", "status": "PAUSED", "category": "Website"},
    )
    assert updated.status_code == 200
    assert updated.json()["name"] == "Client Website Refresh"
    assert updated.json()["status"] == "PAUSED"

    archived = client.post(f"/api/projects/{project['id']}/archive")
    assert archived.status_code == 200
    assert archived.json()["status"] == "ARCHIVED"
    assert not any(row["id"] == project["id"] for row in client.get("/api/projects").json())
    assert any(row["id"] == project["id"] for row in client.get("/api/projects?include_archived=true").json())


def test_project_lists_are_scoped_by_workspace(client):
    ride_project = client.post("/api/projects", json={"name": "Boat Cruise"}).json()
    freelance = client.post("/api/workspaces", json={"name": "Freelance Work"}).json()
    freelance_project = client.post("/api/projects", json={"workspace_id": freelance["id"], "name": "Client Website"}).json()

    active_projects = client.get("/api/projects").json()
    assert [row["id"] for row in active_projects] == [ride_project["id"]]

    client.post(f"/api/workspaces/{freelance['id']}/select")
    switched_projects = client.get("/api/projects").json()
    assert [row["id"] for row in switched_projects] == [freelance_project["id"]]


def test_work_item_project_assignment_validation_and_unassign(client):
    ride_project = client.post("/api/projects", json={"name": "Boat Cruise"}).json()
    created = client.post("/api/work-logs", json={"raw_text": "Reviewed booking copy."}).json()["extracted_items"][0]
    assert created["project_id"] is None

    assigned = client.put(f"/api/work-items/{created['id']}", json={"project_id": ride_project["id"]})
    assert assigned.status_code == 200
    assert assigned.json()["project_id"] == ride_project["id"]
    assert assigned.json()["project_name"] == "Boat Cruise"

    unassigned = client.put(f"/api/work-items/{created['id']}", json={"project_id": None})
    assert unassigned.status_code == 200
    assert unassigned.json()["project_id"] is None

    freelance = client.post("/api/workspaces", json={"name": "Freelance Work"}).json()
    freelance_project = client.post("/api/projects", json={"workspace_id": freelance["id"], "name": "Client Website"}).json()
    bad_assignment = client.put(f"/api/work-items/{created['id']}", json={"project_id": freelance_project["id"]})
    assert bad_assignment.status_code == 400

    archived = client.post(f"/api/projects/{ride_project['id']}/archive").json()
    assert archived["status"] == "ARCHIVED"
    archived_assignment = client.put(f"/api/work-items/{created['id']}", json={"project_id": ride_project["id"]})
    assert archived_assignment.status_code == 400


def test_project_filtered_reports_and_timeline(client):
    project = client.post("/api/projects", json={"name": "Client Website"}).json()
    included = client.post(
        "/api/work-logs",
        json={"raw_text": "Updated the landing-page copy.", "project_id": project["id"]},
    ).json()["extracted_items"][0]
    excluded = client.post("/api/work-logs", json={"raw_text": "Prepared an unrelated invoice."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{included['id']}/confirm")
    client.post(f"/api/work-items/{excluded['id']}/confirm")

    timeline = client.get(f"/api/timeline?project_id={project['id']}").json()
    timeline_ids = {item["id"] for group in timeline for item in group["items"]}
    assert timeline_ids == {included["id"]}

    preview = client.post(
        "/api/reports/preview",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31", "project_id": project["id"]},
    ).json()
    assert [item["id"] for item in preview["confirmed"]] == [included["id"]]

    report = client.post(
        "/api/reports",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31", "project_id": project["id"]},
    ).json()
    assert report["project_id"] == project["id"]
    assert "invoice" not in report["draft_markdown"].lower()


def test_report_prompt_uses_selected_workspace_and_audience(client, monkeypatch):
    client.put(
        "/api/workspaces/default",
        json={"name": "Freelance Work", "user_name": "Alfy", "report_audience": "Client"},
    )
    client.put("/api/settings", json={"key": "selected_model", "value": "mock-model"})
    created = client.post(
        "/api/work-logs",
        json={"raw_text": "Met the client, reviewed the draft poster, recorded wording changes, and sent the revised design for approval."},
    ).json()["extracted_items"][0]
    client.post(f"/api/work-items/{created['id']}/confirm")

    captured = {}
    monkeypatch.setattr(OllamaProvider, "health_check", lambda self: {"available": True, "models": ["mock-model"], "message": "Ollama is reachable."})

    def fake_generate_text(self, model, prompt, system="", options=None):
        captured["prompt"] = prompt
        return "# Weekly Work Report\n\nReport for Freelance Work.\n"

    monkeypatch.setattr(OllamaProvider, "generate_text", fake_generate_text)

    report = client.post(
        "/api/reports",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    )
    assert report.status_code == 200
    assert "Freelance Work" in captured["prompt"]
    assert "Audience: Client" in captured["prompt"]
    assert "Ride Yanga" not in captured["prompt"]
    assert "CTO" not in captured["prompt"]


def test_generic_report_labels_keep_legacy_report_types_usable(client):
    report_types = client.get("/api/reports/types").json()
    assert "Monthly Work Summary" in report_types
    assert "Stakeholder Update" in report_types
    assert "Monthly Engineering Summary" not in report_types
    assert "Boss Progress Update" not in report_types

    created = client.post("/api/work-logs", json={"raw_text": "Prepared a weekly work summary."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{created['id']}/confirm")
    report = client.post(
        "/api/reports",
        json={"report_type": "Monthly Engineering Summary", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    )
    assert report.status_code == 200
    body = report.json()
    assert body["report_type"] == "Monthly Engineering Summary"
    assert body["title"].startswith("Monthly Work Summary")
    assert any(row["report_type"] == "Monthly Engineering Summary" for row in client.get("/api/reports").json())


def test_ollama_semantic_extraction_creates_multiple_review_items(client, monkeypatch):
    client.put("/api/settings", json={"key": "selected_model", "value": "mock-model"})

    monkeypatch.setattr(OllamaProvider, "health_check", lambda self: {"available": True, "models": ["mock-model"], "message": "Ollama is reachable."})

    def fake_generate_structured(self, model, prompt, schema, system=""):
        assert model == "mock-model"
        assert schema in {SemanticExtractionSchema, RawSemanticExtractionSchema}
        return SemanticExtractionSchema(
            events=[
                SemanticEventSchema(
                    event_subject="Boat Cruise category images",
                    initial_context=["The boat cruise category icons were being reviewed."],
                    actions_performed=[SemanticFactSchema(fact_type="ACTION", subject="Boat Cruise category images", statement="Fixed the boat cruise category images after finding one icon was wrong.", certainty="CONFIRMED", status="COMPLETED")],
                    findings=[SemanticFactSchema(fact_type="FINDING", subject="Boat Cruise category images", statement="One of the category icons was wrong.", certainty="CONFIRMED", status="COMPLETED")],
                    current_status="COMPLETED",
                ),
                SemanticEventSchema(
                    event_subject="Company-wide notification capability",
                    initial_context=["The notification requirement was clarified as company announcements."],
                    findings=[SemanticFactSchema(fact_type="FINDING", subject="Dashboard notification configuration", statement="Notification toggles exist in the dashboard.", certainty="CONFIRMED", status="COMPLETED")],
                    open_questions=[SemanticFactSchema(fact_type="OPEN_QUESTION", subject="Broadcast sending", statement="Whether broadcast sending is implemented is still unconfirmed.", certainty="UNVERIFIED", status="OPEN")],
                    pending_actions=[SemanticFactSchema(fact_type="PENDING_ACTION", subject="Broadcast sending", statement="Confirm whether broadcast sending is actually implemented.", certainty="UNVERIFIED", status="PENDING")],
                    current_status="ONGOING",
                ),
            ]
        )

    monkeypatch.setattr(OllamaProvider, "generate_structured", fake_generate_structured)

    response = client.post("/api/work-logs", json={"raw_text": MULTI_EVENT_NOTE})
    assert response.status_code == 200
    body = response.json()
    assert body["analysis_mode"] == "OLLAMA"
    assert body["analysis_model"] == "mock-model"
    assert len(body["extracted_items"]) == 2
    titles = " ".join(item["title"].lower() for item in body["extracted_items"])
    assert "boat cruise" in titles
    assert "notification" in titles


def test_invalid_ai_semantic_output_falls_back_safely(client, monkeypatch):
    client.put("/api/settings", json={"key": "selected_model", "value": "mock-model"})

    monkeypatch.setattr(OllamaProvider, "health_check", lambda self: {"available": True, "models": ["mock-model"], "message": "Ollama is reachable."})
    monkeypatch.setattr(OllamaProvider, "generate_structured", lambda *args, **kwargs: (_ for _ in ()).throw(ValueError("invalid json")))

    response = client.post("/api/work-logs", json={"raw_text": NOTIFICATION_NOTE})
    assert response.status_code == 200
    body = response.json()
    assert body["analysis_mode"] == "EVIDENCE_ONLY"
    assert body["extracted_items"]


def test_semantic_validation_demotes_unsupported_confirmation():
    raw_text = NOTIFICATION_NOTE
    result = SemanticExtractionSchema(
        events=[
            SemanticEventSchema(
                event_subject="Company-wide notification capability",
                findings=[
                    SemanticFactSchema(
                        fact_type="FINDING",
                        subject="Company-wide notification capability",
                        statement="Company notifications were activated from the dashboard.",
                        certainty="CONFIRMED",
                        status="COMPLETED",
                    )
                ],
            )
        ]
    )
    normalized = normalize_semantic_result(result, raw_text, [])
    event = normalized.events[0]
    assert event.current_status == "ONGOING"
    assert any(fact.fact_type == "PENDING_ACTION" for fact in event.pending_actions)
    assert any("confirm" in fact.statement.lower() for fact in event.pending_actions)


def test_multiple_topic_note_splits_into_multiple_review_items(client):
    response = client.post("/api/work-logs", json={"raw_text": MULTI_EVENT_NOTE})
    assert response.status_code == 200
    body = response.json()
    assert len(body["extracted_items"]) >= 2
    titles = " ".join(item["title"].lower() for item in body["extracted_items"])
    assert "boat cruise" in titles
    assert "notification" in titles


def test_report_preview_filters_confirmed_and_inferred(client):
    first = client.post("/api/work-logs", json={"raw_text": "Implemented a booking validation fix."}).json()["extracted_items"][0]
    second = client.post("/api/work-logs", json={"raw_text": "Investigated notification toggle behavior."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{first['id']}/confirm")

    preview = client.post(
        "/api/reports/preview",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    )
    assert preview.status_code == 200
    body = preview.json()
    assert any(item["id"] == first["id"] for item in body["confirmed"])
    assert any(item["id"] == second["id"] for item in body["inferred_needs_review"])


def test_fts_retrieval_finds_work_log(client):
    client.post("/api/work-logs", json={"raw_text": "Resolved permalink issue in WordPress navigation testing."})
    response = client.get("/api/search?q=permalink")
    assert response.status_code == 200
    assert response.json()


def test_repository_validation_rejects_non_git(client, tmp_path):
    response = client.post(
        "/api/repositories",
        json={"name": "Not Git", "local_path": str(tmp_path), "role": "OTHER"},
    )
    assert response.status_code == 400


def test_git_scan_collects_commit_metadata(client, tmp_path):
    if not shutil.which("git"):
        return
    repo = tmp_path / "repo"
    repo.mkdir()
    subprocess.run(["git", "init"], cwd=repo, check=True, capture_output=True)
    subprocess.run(["git", "config", "user.email", "test@example.com"], cwd=repo, check=True)
    subprocess.run(["git", "config", "user.name", "Tester"], cwd=repo, check=True)
    (repo / "notifications.txt").write_text("promo notifications", encoding="utf-8")
    subprocess.run(["git", "add", "."], cwd=repo, check=True)
    subprocess.run(["git", "commit", "-m", "Add promo notification notes"], cwd=repo, check=True, capture_output=True)

    created = client.post("/api/repositories", json={"name": "Working Repo", "local_path": str(repo), "role": "WORKING_SANDBOX"})
    assert created.status_code == 200
    scan = client.post(f"/api/git/scan-now/{created.json()['id']}")
    assert scan.status_code == 200
    assert scan.json()["new_commits"] == 1

    review = client.get("/api/work-items?status=REVIEW")
    assert any("promo notification" in f"{item['title']} {item['summary']}".lower() for item in review.json())
    evidence = client.get("/api/search?q=notifications")
    assert evidence.status_code == 200


def test_sandbox_to_canonical_detection_works_when_canonical_scanned_second(client, tmp_path):
    if not shutil.which("git"):
        return
    sandbox = make_repo(tmp_path / "working", "Add company notification workflow", "src/notifications/company.txt")
    canonical = make_repo(tmp_path / "dashboard", "Add company notification workflow", "src/notifications/company.txt")

    canonical_repo = client.post(
        "/api/repositories",
        json={"name": "Dashboard API", "local_path": str(canonical), "role": "DASHBOARD_API"},
    ).json()
    sandbox_repo = client.post(
        "/api/repositories",
        json={
            "name": "Working Repo",
            "local_path": str(sandbox),
            "role": "WORKING_SANDBOX",
            "promotes_to_repository_id": canonical_repo["id"],
        },
    ).json()

    assert client.post(f"/api/git/scan-now/{sandbox_repo['id']}").status_code == 200
    assert client.post(f"/api/git/scan-now/{canonical_repo['id']}").status_code == 200
    relationships = client.get("/api/evidence/relationships").json()
    assert any(row["relationship_type"] == "PROMOTED_TO" for row in relationships)


def test_chat_this_week_uses_week_date_filter(client):
    old = client.post("/api/work-logs", json={"raw_text": "Implemented old booking report flow."}).json()["extracted_items"][0]
    client.put(f"/api/work-items/{old['id']}", json={"work_date": "2020-01-01", "status": "CONFIRMED"})
    current = client.post("/api/work-logs", json={"raw_text": "Investigated current week notification flow."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{current['id']}/confirm")

    response = client.post("/api/chat", json={"message": "What did I do this week?"}).json()
    answer = response["answer"]
    assert "current week" in answer.lower()
    assert "old booking" not in answer.lower()
    assert "work item" in response["evidence_summary"].lower()


def test_docx_import_stores_document_and_review_items(client, tmp_path):
    doc = Document()
    doc.add_paragraph("Weekly report: investigated dashboard notifications and fixed a booking issue.")
    path = tmp_path / "report.docx"
    doc.save(path)
    with path.open("rb") as handle:
        response = client.post("/api/imports", files={"files": ("report.docx", handle, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")})
    assert response.status_code == 200
    body = response.json()[0]
    assert body["duplicate"] is False
    assert body["extracted_items"]


def test_report_excludes_promoted_sandbox_duplicate(client, tmp_path):
    if not shutil.which("git"):
        return
    sandbox = make_repo(tmp_path / "working2", "Add booking confirmation validation", "src/booking/confirm.txt")
    canonical = make_repo(tmp_path / "dashboard2", "Add booking confirmation validation", "src/booking/confirm.txt")
    canonical_repo = client.post("/api/repositories", json={"name": "Dashboard API 2", "local_path": str(canonical), "role": "DASHBOARD_API"}).json()
    sandbox_repo = client.post(
        "/api/repositories",
        json={"name": "Working Repo 2", "local_path": str(sandbox), "role": "WORKING_SANDBOX", "promotes_to_repository_id": canonical_repo["id"]},
    ).json()
    client.post(f"/api/git/scan-now/{sandbox_repo['id']}")
    client.post(f"/api/git/scan-now/{canonical_repo['id']}")
    for item in client.get("/api/work-items?status=REVIEW").json():
        if "booking confirmation" in f"{item['title']} {item['summary']}".lower():
            client.post(f"/api/work-items/{item['id']}/confirm")

    report = client.post(
        "/api/reports",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    ).json()
    work_bullets = [line for line in report["draft_markdown"].lower().splitlines() if line.startswith("- ") and "booking" in line]
    assert len(work_bullets) <= 1


def test_notification_investigation_report_is_conservative_and_keeps_pending_work(client):
    created = client.post("/api/work-logs", json={"raw_text": NOTIFICATION_NOTE}).json()["extracted_items"]
    for item in created:
        client.post(f"/api/work-items/{item['id']}/confirm")

    report = client.post(
        "/api/reports",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    ).json()
    draft = report["draft_markdown"].lower()
    forbidden = [
        "company notifications were activated",
        "company notification capability was confirmed",
        "broadcast notifications were implemented",
        "no pending work",
        "marketing initiatives",
        "monitoring the broadcast flow for effectiveness",
    ]
    for phrase in forbidden:
        assert phrase not in draft
    assert "investigat" in draft
    assert "notification" in draft
    assert "open questions / pending work" in draft
    assert "confirm" in draft
    assert "broadcast" in draft
    assert "okay so today" not in draft


def test_atomic_decomposition_splits_confirmed_finding_open_question_and_pending_action(client):
    created = client.post("/api/work-logs", json={"raw_text": NOTIFICATION_NOTE}).json()["extracted_items"]
    for item in created:
        client.post(f"/api/work-items/{item['id']}/confirm")
    items = client.get("/api/work-items?status=CONFIRMED").json()
    notification_items = [
        item
        for item in items
        if "notification" in f"{item['title']} {item['summary']} {item.get('findings') or ''} {item.get('pending_work') or ''}".lower()
    ]
    assert notification_items
    notification_item = next(
        row
        for row in notification_items
        if row.get("findings") and row.get("pending_work")
    )

    class Obj:
        def __init__(self, data):
            self.__dict__.update(data)

    plan = build_report_plan("Weekly Work Report", "2026-07-06", "2026-07-12", [Obj(notification_item)])
    fact_types = {fact.fact_type for fact in plan.facts}
    assert {"FINDING", "OPEN_QUESTION", "PENDING_ACTION"}.issubset(fact_types)
    finding = next(fact for fact in plan.facts if fact.fact_type == "FINDING")
    open_question = next(fact for fact in plan.facts if fact.fact_type == "OPEN_QUESTION")
    pending = next(fact for fact in plan.facts if fact.fact_type == "PENDING_ACTION")
    assert "toggles" in finding.statement.lower()
    assert finding.certainty == "CONFIRMED"
    assert "has not yet been confirmed" in open_question.statement.lower()
    assert pending.statement.lower().startswith("inspect") or pending.statement.lower().startswith("verify")


def test_report_sections_are_professional_distinct_and_non_redundant(client):
    created = client.post("/api/work-logs", json={"raw_text": NOTIFICATION_NOTE}).json()["extracted_items"]
    for item in created:
        client.put(f"/api/work-items/{item['id']}", json={"work_date": "2026-07-08"})
        client.post(f"/api/work-items/{item['id']}/confirm")
    report = client.post(
        "/api/reports",
        json={"report_type": "Weekly Work Report", "date_from": "2026-07-06", "date_to": "2026-07-12"},
    ).json()
    draft = report["draft_markdown"]
    lower = draft.lower()
    assert "okay so today" not in lower
    assert "notification-related configuration toggles were identified in the dashboard" in lower

    findings = section_lines(draft, "Confirmed Findings")
    pending = section_lines(draft, "Open Questions / Pending Work")
    next_steps = section_lines(draft, "Next Steps")
    assert findings
    assert pending
    assert next_steps
    assert all(not any(word in line.lower() for word in ["need to confirm", "whether", "unconfirmed", "not yet"]) for line in findings)
    assert normalized_semantic_text(pending[0]) != normalized_semantic_text(next_steps[0])
    assert any(line.lower().startswith(("inspect", "verify")) for line in next_steps)


def test_report_refinement_creates_revision_and_restore(client):
    item = client.post("/api/work-logs", json={"raw_text": "Investigated dashboard notification toggles and still need to confirm broadcast flow."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{item['id']}/confirm")
    report = client.post(
        "/api/reports",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    ).json()
    refined = client.post(f"/api/reports/{report['id']}/refine", json={"instruction": "Do not say I implemented notifications. I only investigated them."})
    assert refined.status_code == 200
    revisions = client.get(f"/api/reports/{report['id']}/revisions").json()
    assert len(revisions) >= 2
    restored = client.post(f"/api/reports/{report['id']}/revisions/{revisions[0]['id']}/restore")
    assert restored.status_code == 200


def test_docx_export_renders_markdown_without_literal_syntax(client, tmp_path):
    item = client.post("/api/work-logs", json={"raw_text": "Investigated notification toggles and still need to confirm broadcast flow."}).json()["extracted_items"][0]
    client.post(f"/api/work-items/{item['id']}/confirm")
    report = client.post(
        "/api/reports",
        json={"report_type": "Weekly Work Report", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    ).json()
    response = client.get(f"/api/reports/{report['id']}/export/docx")
    assert response.status_code == 200
    path = tmp_path / "report.docx"
    path.write_bytes(response.content)
    doc = Document(path)
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    assert "**" not in text
    assert "###" not in text
    assert "Audience:" not in text
    assert "Tone:" not in text
    assert "Thank you for your attention" not in text


def test_pptx_export_uses_structured_slides_not_raw_report_dump(client, tmp_path):
    created = client.post("/api/work-logs", json={"raw_text": NOTIFICATION_NOTE}).json()["extracted_items"]
    for item in created:
        client.post(f"/api/work-items/{item['id']}/confirm")
    report = client.post(
        "/api/reports",
        json={"report_type": "Meeting Presentation", "date_from": "2000-01-01", "date_to": "2999-12-31"},
    ).json()
    response = client.get(f"/api/reports/{report['id']}/export/pptx")
    assert response.status_code == 200
    path = tmp_path / "report.pptx"
    path.write_bytes(response.content)
    prs = Presentation(path)
    assert len(prs.slides) >= 3
    all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "Audience" not in all_text
    assert "Tone" not in all_text
    assert "**" not in all_text
    assert "okay so today" not in all_text.lower()
    assert "Notification-related configuration toggles" in all_text
    assert "has not yet been confirmed" in all_text


def test_presentation_plan_consumes_atomic_facts_not_raw_report_paragraphs(client):
    created = client.post("/api/work-logs", json={"raw_text": NOTIFICATION_NOTE}).json()["extracted_items"]
    for item in created:
        client.post(f"/api/work-items/{item['id']}/confirm")
    confirmed_items = client.get("/api/work-items?status=CONFIRMED").json()
    item = next(
        row
        for row in confirmed_items
        if "notification" in f"{row['title']} {row['summary']} {row.get('findings') or ''} {row.get('pending_work') or ''}".lower()
        and row.get("findings")
        and row.get("pending_work")
    )

    class Obj:
        def __init__(self, data):
            self.__dict__.update(data)

    plan = build_report_plan("Meeting Presentation", "2026-07-06", "2026-07-12", [Obj(item)])
    presentation = presentation_plan_from_report_plan(plan)
    bullets = "\n".join(bullet for slide in presentation["slides"] for bullet in slide["bullets"])
    assert "okay so today" not in bullets.lower()
    assert "Notification-related configuration toggles were identified in the dashboard." in bullets
    assert "The existence and operation of a company-wide broadcast notification flow has not yet been confirmed." in bullets


def make_repo(path: Path, message: str, file_path: str) -> Path:
    path.mkdir()
    subprocess.run(["git", "init"], cwd=path, check=True, capture_output=True)
    subprocess.run(["git", "config", "user.email", "test@example.com"], cwd=path, check=True)
    subprocess.run(["git", "config", "user.name", "Tester"], cwd=path, check=True)
    target = path / file_path
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(message, encoding="utf-8")
    subprocess.run(["git", "add", "."], cwd=path, check=True)
    subprocess.run(["git", "commit", "-m", message], cwd=path, check=True, capture_output=True)
    return path


def section_lines(markdown: str, heading: str) -> list[str]:
    lines = markdown.splitlines()
    result = []
    in_section = False
    for line in lines:
        if line == f"## {heading}":
            in_section = True
            continue
        if in_section and line.startswith("## "):
            break
        if in_section and line.startswith("- "):
            result.append(line[2:])
    return result
